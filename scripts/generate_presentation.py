"""Generates docs/presentation/Migration_Wizard_Praktikum_Presentation.pptx.

One-shot content generator, not part of the app itself. Re-run after editing
this file to regenerate the deck. Screenshots in docs/presentation/assets/
are real renders of the actual app/report (see git history for how they
were captured) — keep them in sync if the UI changes meaningfully.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

ASSETS = Path(__file__).resolve().parent.parent / "docs" / "presentation" / "assets"
OUT_PATH = Path(__file__).resolve().parent.parent / "docs" / "presentation" / "Migration_Wizard_Praktikum_Presentation.pptx"

NAVY = RGBColor(0x1B, 0x1E, 0x28)
BLUE = RGBColor(0x3F, 0x6F, 0xE0)
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x86)
LIGHT_BLUE_BG = RGBColor(0xDC, 0xE6, 0xFF)
GREEN = RGBColor(0x1B, 0x5E, 0x20)
GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)
AMBER = RGBColor(0xE6, 0x51, 0x00)
AMBER_BG = RGBColor(0xFF, 0xF3, 0xE0)
RED = RGBColor(0xB7, 0x1C, 0x1C)
GREY = RGBColor(0x54, 0x6E, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF7, 0xFB, 0xFE)
LINE = RGBColor(0xC3, 0xD6, 0xE0)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _send_to_back(slide, shape):
    spTree = slide.shapes._spTree
    spTree.remove(shape._element)
    spTree.insert(2, shape._element)


def blank_slide(prs, fill=BG):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = fill
    bg.line.fill.background()
    bg.shadow.inherit = False
    _send_to_back(slide, bg)
    return slide


def add_text(slide, left, top, width, height, text, size=18, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, font="Segoe UI", anchor=None, italic=False, line_spacing=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = font
    return box


def add_card(slide, left, top, width, height, fill=WHITE, line_color=LINE, radius=True, shadow=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    card = slide.shapes.add_shape(shape_type, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = line_color
    card.line.width = Pt(1)
    card.shadow.inherit = shadow
    if radius:
        try:
            card.adjustments[0] = 0.07
        except Exception:
            pass
    return card


def add_pill(slide, left, top, width, height, text, fg, bg, size=12):
    pill = add_card(slide, left, top, width, height, fill=bg, line_color=bg)
    try:
        pill.adjustments[0] = 0.5
    except Exception:
        pass
    tf = pill.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = fg
    run.font.name = "Segoe UI"
    return pill


def add_icon_circle(slide, cx, cy, diameter, glyph, fg=WHITE, bg=BLUE, size=28):
    left = cx - diameter / 2
    top = cy - diameter / 2
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    circ.fill.solid()
    circ.fill.fore_color.rgb = bg
    circ.line.fill.background()
    circ.shadow.inherit = False
    tf = circ.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = glyph
    run.font.size = Pt(size)
    run.font.color.rgb = fg
    run.font.bold = True
    return circ


def add_arrow(slide, x1, y1, x2, y2, color=BLUE, weight=2.25):
    """A simple horizontal right-pointing arrow autoshape between two points.

    A connector + manually-injected <a:tailEnd> XML looked right in
    python-pptx (non-validating) but produced OOXML PowerPoint itself
    refused to open — child-element ordering in <a:ln> is schema-strict.
    A plain autoshape avoids hand-written XML entirely.
    """
    width = abs(x2 - x1)
    height = Pt(14)
    top = min(y1, y2) - height / 2
    left = min(x1, x2)
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    arrow.shadow.inherit = False
    return arrow


def add_header(slide, kicker, title, kicker_color=BLUE):
    add_text(slide, Inches(0.6), Inches(0.3), Inches(11), Inches(0.35), kicker.upper(),
              size=12, bold=True, color=kicker_color)
    add_text(slide, Inches(0.6), Inches(0.62), Inches(12.0), Inches(0.75), title,
              size=27, bold=True, color=NAVY)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.32), Inches(1.1), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()
    line.shadow.inherit = False


def add_picture_framed(slide, img_path, left, top, max_w, max_h, caption=None, border=True):
    from PIL import Image
    with Image.open(img_path) as im:
        iw, ih = im.size
    scale = min(max_w / iw, max_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    frame_pad = Pt(4)
    if border:
        frame = add_card(slide, left - Pt(6), top - Pt(6), w + Pt(12), h + Pt(12), fill=WHITE, line_color=LINE, radius=True, shadow=True)
    pic = slide.shapes.add_picture(str(img_path), left, top, width=w, height=h)
    if caption:
        add_text(slide, left, top + h + Pt(8), max_w, Inches(0.4), caption, size=12.5, italic=True, color=GREY,
                  align=PP_ALIGN.CENTER)
    return pic, w, h


def icon_label_row(slide, items, top, total_width=Inches(11.5), left=Inches(0.9), circle_d=Inches(0.85),
                    icon_size=30, label_size=14, sub_size=11.5, fg=WHITE, bg=BLUE):
    n = len(items)
    slot = total_width / n
    for i, (glyph, label, sub) in enumerate(items):
        cx = left + slot * i + slot / 2
        add_icon_circle(slide, cx, top + circle_d / 2, circle_d, glyph, fg=fg, bg=bg, size=icon_size)
        add_text(slide, cx - slot / 2 + Inches(0.1), top + circle_d + Pt(8), slot - Inches(0.2), Inches(0.4),
                  label, size=label_size, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        if sub:
            add_text(slide, cx - slot / 2 + Inches(0.1), top + circle_d + Pt(8) + Inches(0.42), slot - Inches(0.2), Inches(0.7),
                      sub, size=sub_size, color=GREY, align=PP_ALIGN.CENTER, line_spacing=1.15)


def set_table_style(table, header_bg=DARK_BLUE, header_fg=WHITE, body_size=13, header_size=13):
    for r_idx, row in enumerate(table.rows):
        for c_idx, cell in enumerate(row.cells):
            cell.margin_left = Pt(8)
            cell.margin_right = Pt(8)
            cell.margin_top = Pt(4)
            cell.margin_bottom = Pt(4)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_bg if r_idx == 0 else (WHITE if r_idx % 2 else RGBColor(0xEE, 0xF7, 0xFB))
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if (r_idx == 0 or c_idx > 0) else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.size = Pt(header_size if r_idx == 0 else body_size)
                    run.font.color.rgb = header_fg if r_idx == 0 else NAVY
                    run.font.bold = r_idx == 0
                    run.font.name = "Segoe UI"


def fill_table(slide, left, top, width, height, rows_data):
    rows, cols = len(rows_data), len(rows_data[0])
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    for r, row_vals in enumerate(rows_data):
        for c, val in enumerate(row_vals):
            table.cell(r, c).text = str(val)
    set_table_style(table)
    return table


def add_bar_chart(slide, left, top, width, height, categories, series_name, values,
                   bar_color=BLUE, title=None, horizontal=False):
    """A real native PowerPoint chart (editable in PowerPoint), not a picture."""
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(series_name, values)
    chart_type = XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED
    graphic_frame = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data)
    chart = graphic_frame.chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(13)
    plot.data_labels.font.bold = True
    plot.data_labels.font.color.rgb = NAVY
    if not horizontal:
        plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = bar_color
    series.format.line.fill.background()
    chart.category_axis.format.line.color.rgb = LINE
    chart.category_axis.tick_labels.font.size = Pt(13)
    chart.category_axis.tick_labels.font.color.rgb = NAVY
    chart.value_axis.visible = False
    chart.value_axis.has_major_gridlines = False
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
        chart.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
    else:
        chart.has_title = False
    return chart


# ─────────────────────────────────────────────────────────────────────────────

def build() -> Presentation:
    prs = new_presentation()

    # ── 1. Title ─────────────────────────────────────────────────────────────
    s = blank_slide(prs, fill=NAVY)
    add_icon_circle(s, Inches(2.3), Inches(1.5), Inches(1.0), "\U0001FA9F", fg=WHITE, bg=DARK_BLUE, size=36)  # window
    add_text(s, Inches(2.85), Inches(1.05), Inches(0.9), Inches(0.9), "→", size=36, bold=True, color=BLUE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_icon_circle(s, Inches(4.0), Inches(1.5), Inches(1.0), "\U0001F427", fg=WHITE, bg=GREEN, size=36)  # penguin
    add_text(s, Inches(0.7), Inches(2.4), Inches(11.5), Inches(1.3), "Migration Wizard", size=54, bold=True, color=WHITE)
    add_text(s, Inches(0.7), Inches(3.55), Inches(11.5), Inches(0.9),
              "Windows 11  →  Linux Mint, automated", size=22, color=LIGHT_BLUE_BG)
    add_text(s, Inches(0.7), Inches(6.5), Inches(8), Inches(0.5),
              "Japhet Kofi Appau Arthur  ·  Mobile Computing Seminar · Praktikum", size=15, color=RGBColor(0xAE, 0xC2, 0xE8))

    # ── 1b. Building on prior work ────────────────────────────────────────────
    s = blank_slide(prs)
    add_header(s, "Where this builds from", "Extending a term paper into a working tool")
    add_card(s, Inches(0.8), Inches(1.85), Inches(11.7), Inches(1.5), fill=WHITE)
    add_text(s, Inches(1.1), Inches(2.0), Inches(11.1), Inches(0.5), "TERM PAPER", size=12.5, bold=True, color=BLUE)
    add_text(s, Inches(1.1), Inches(2.35), Inches(11.1), Inches(0.55),
              "“Digital Sovereignty Through Semi-Automated Migration from Windows to Linux”", size=18, bold=True, color=NAVY)
    add_text(s, Inches(1.1), Inches(2.92), Inches(11.1), Inches(0.4),
              "Japhet Arthur — Julius-Maximilians-Universität Würzburg, Lehrstuhl für Kommunikationsnetze", size=13.5, italic=True, color=GREY)
    add_text(s, Inches(0.8), Inches(3.65), Inches(11.7), Inches(0.6), "RESEARCH QUESTION IT ASKED", size=13, bold=True, color=DARK_BLUE)
    add_text(s, Inches(0.8), Inches(4.0), Inches(11.7), Inches(0.85),
              "How can semi-automated migration frameworks reduce technical barriers to OS-level\nsovereignty restoration while maintaining user control, data integrity, and workflow continuity?",
              size=16, color=NAVY, line_spacing=1.25)
    add_card(s, Inches(0.8), Inches(5.05), Inches(11.7), Inches(1.6), fill=AMBER_BG, line_color=AMBER_BG)
    add_text(s, Inches(1.1), Inches(5.25), Inches(11.1), Inches(0.4), "WHAT IT BUILT — AND STATED AS LIMITATIONS", size=13, bold=True, color=AMBER)
    add_text(s, Inches(1.1), Inches(5.6), Inches(11.1), Inches(0.9),
              "A working 5-phase pipeline (Assess → Review → Prepare → Migrate → Verify) — but compatibility\nmapping was “heuristic... not guaranteeing feature parity,” with “enhanced compatibility detection” named as future work.",
              size=14.5, color=NAVY, line_spacing=1.25)

    # ── 2. Motivation ─────────────────────────────────────────────────────────
    s = blank_slide(prs)
    add_header(s, "Motivation", "Three things to move. Zero tools that move all three.")
    icon_label_row(s, [
        ("\U0001F4C1", "Files", "Documents, photos,\nmusic, ..."),
        ("\U0001F5A5", "Apps", "Find the right\nLinux equivalent"),
        ("⚙", "Settings", "Wallpaper, theme,\nshortcuts"),
    ], top=Inches(2.0), bg=BLUE)
    add_card(s, Inches(1.5), Inches(4.5), Inches(10.3), Inches(1.6), fill=WHITE)
    add_text(s, Inches(1.5), Inches(4.75), Inches(10.3), Inches(1.1),
              "Doing this by hand means researching dozens of app equivalents,\nmanually copying files, and rebuilding your desktop from scratch.",
              size=18, color=NAVY, align=PP_ALIGN.CENTER, line_spacing=1.3)

    # ── 2b. Digital Sovereignty ────────────────────────────────────────────────
    s = blank_slide(prs)
    add_header(s, "Why it matters", "This is a digital sovereignty problem")
    add_text(s, Inches(0.9), Inches(1.85), Inches(11.0), Inches(0.9),
              "“The capacity to exercise meaningful control over digital infrastructure, data flows,\nand computational processes” — autonomy, control, transparency.",
              size=17, color=NAVY, line_spacing=1.3)
    add_text(s, Inches(0.9), Inches(2.55), Inches(11.0), Inches(0.35),
              "Floridi (2020); Pohle & Thiel (2020); Roberts et al. (2021)", size=12.5, italic=True, color=GREY)
    pairs = [
        ("\U0001F510", "Closed → Open", "Proprietary OS/apps → an\nopen-source stack you control"),
        ("☁", "Cloud-tied → Local-first", "Files stay on your machine,\nnot a vendor's cloud"),
        ("\U0001F517", "Locked-in → Portable", "MIT-licensed tool; your data\nleaves in standard formats"),
    ]
    icon_label_row(s, pairs, top=Inches(2.95), circle_d=Inches(1.0), icon_size=28, label_size=15, sub_size=13, bg=DARK_BLUE)
    add_card(s, Inches(1.5), Inches(5.4), Inches(10.3), Inches(1.3), fill=GREEN_BG, line_color=GREEN_BG)
    add_text(s, Inches(1.8), Inches(5.6), Inches(9.7), Inches(0.5), "MEASURED, NOT JUST CLAIMED", size=13, bold=True, color=GREEN)
    add_text(s, Inches(1.8), Inches(5.95), Inches(9.7), Inches(0.6),
              "sovereignty_score = integrity_score + openness_bonus  — every restore reports one",
              size=15, color=NAVY)

    # ── 2c. Preference-behavior gap ─────────────────────────────────────────────
    s = blank_slide(prs)
    add_header(s, "Why it matters", "Europeans want sovereignty, but don't have it yet")
    add_bar_chart(
        s, Inches(0.8), Inches(1.9), Inches(7.6), Inches(4.4),
        categories=["Rate open-source\ncritical for\nsovereignty", "Cite reducing\nUS-vendor\ndependency", "Still run\nWindows in\nEurope"],
        series_name="Share of respondents / market",
        values=[63.2, 47.4, 67.75],
        bar_color=DARK_BLUE,
        title=None,
    )
    add_card(s, Inches(8.7), Inches(1.95), Inches(3.8), Inches(4.3), fill=AMBER_BG, line_color=AMBER_BG)
    add_text(s, Inches(9.0), Inches(2.15), Inches(3.2), Inches(0.4), "THE GAP", size=13, bold=True, color=AMBER)
    add_text(s, Inches(9.0), Inches(2.55), Inches(3.2), Inches(2.0),
              "Stated preference for\nsovereignty far exceeds\nactual behavior — the gap\nis migration barriers,\nnot lack of will.",
              size=14.5, color=NAVY, line_spacing=1.3)
    add_text(s, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5),
              "Wire, State of Digital Sovereignty in Europe (2025) · StatCounter, Desktop OS Market Share Europe (2025)",
              size=12, italic=True, color=GREY)

    # ── 2d. Telemetry & Big Tech ───────────────────────────────────────────────
    s = blank_slide(prs)
    add_header(s, "What users are escaping", "Telemetry is the concrete problem")
    add_text(s, Inches(0.9), Inches(1.8), Inches(11.4), Inches(0.6),
              "“Windows 11 integrates mandatory telemetry that cannot be fully disabled... users must trust\nvendor claims without ability to verify or constrain collection practices.”",
              size=15.5, color=NAVY, line_spacing=1.25, italic=True)
    add_card(s, Inches(0.7), Inches(2.6), Inches(5.7), Inches(3.9), fill=AMBER_BG, line_color=AMBER_BG)
    add_text(s, Inches(1.0), Inches(2.8), Inches(5.1), Inches(0.4), "TYPICAL BIG TECH DEFAULT", size=13, bold=True, color=AMBER)
    typical = ["Diagnostic data sent by default", "Advertising ID per device", "Account + cloud sign-in required",
               "Telemetry scope set by the vendor", "Closed-source OS internals"]
    y = Inches(3.3)
    for t in typical:
        add_text(s, Inches(1.0), y, Inches(5.1), Inches(0.5), "•  " + t, size=14.5, color=NAVY)
        y += Inches(0.62)
    add_card(s, Inches(6.7), Inches(2.6), Inches(5.9), Inches(3.9), fill=GREEN_BG, line_color=GREEN_BG)
    add_text(s, Inches(7.0), Inches(2.8), Inches(5.3), Inches(0.35), "THIS CODEBASE'S NETWORK CALLS", size=12.5, bold=True, color=GREEN)
    add_text(s, Inches(7.0), Inches(3.1), Inches(5.3), Inches(0.4), "counted directly in source, not claimed", size=11.5, italic=True, color=GREY)
    add_bar_chart(
        s, Inches(6.9), Inches(3.5), Inches(5.5), Inches(2.85),
        categories=["Repology\n(opt-in, Expert\nmode only)", "Telemetry", "Ads /\ntracking", "Remote\nanalytics"],
        series_name="Outbound call sites",
        values=[2, 0, 0, 0],
        bar_color=GREEN,
    )

    # ── 3. Approach ──────────────────────────────────────────────────────────
    s = blank_slide(prs)
    add_header(s, "Approach", "One guided pipeline, two machines")
    add_card(s, Inches(0.6), Inches(1.9), Inches(5.0), Inches(1.6), fill=GREEN_BG, line_color=GREEN_BG)
    add_icon_circle(s, Inches(1.25), Inches(2.7), Inches(0.8), "1", fg=WHITE, bg=GREEN, size=22)
    add_text(s, Inches(1.85), Inches(2.15), Inches(3.6), Inches(0.5), "Prepare (Windows)", size=17, bold=True, color=NAVY)
    add_text(s, Inches(1.85), Inches(2.6), Inches(3.6), Inches(0.8), "Scan → map apps → pack a bundle", size=14, color=GREY)
    add_text(s, Inches(5.9), Inches(2.55), Inches(0.7), Inches(0.7), "→", size=30, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_card(s, Inches(6.7), Inches(1.9), Inches(5.0), Inches(1.6), fill=LIGHT_BLUE_BG, line_color=LIGHT_BLUE_BG)
    add_icon_circle(s, Inches(7.35), Inches(2.7), Inches(0.8), "2", fg=WHITE, bg=DARK_BLUE, size=22)
    add_text(s, Inches(7.95), Inches(2.15), Inches(3.6), Inches(0.5), "Restore (Linux Mint)", size=17, bold=True, color=NAVY)
    add_text(s, Inches(7.95), Inches(2.6), Inches(3.6), Inches(0.8), "Unzip → restore → verify → report", size=14, color=GREY)
    pic, w, h = add_picture_framed(s, ASSETS / "mode_page.png", Inches(2.9), Inches(3.85), Inches(7.5), Inches(3.3))
    add_text(s, Inches(0.6), Inches(3.95), Inches(2.1), Inches(1.0),
              "Same destination,\nyour choice of how\nmuch help along\nthe way →", size=14, italic=True, color=GREY, line_spacing=1.2)

    # ── 4. Architecture ───────────────────────────────────────────────────────
    s = blank_slide(prs)
    add_header(s, "Architecture", "Clean layers, thin Qt window")
    layers = [
        ("inventory/", "collect", GREEN_BG, GREEN),
        ("analysis/", "match", LIGHT_BLUE_BG, DARK_BLUE),
        ("services/", "orchestrate", AMBER_BG, AMBER),
        ("qt_ui/", "present (MVC)", LIGHT_BLUE_BG, DARK_BLUE),
    ]
    x = Inches(0.9)
    w = Inches(2.75)
    for i, (name, verb, bg, fg) in enumerate(layers):
        add_card(s, x, Inches(2.4), w, Inches(1.5), fill=bg, line_color=bg)
        add_text(s, x, Inches(2.62), w, Inches(0.5), name, size=18, bold=True, color=fg, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(3.15), w, Inches(0.5), verb, size=13, color=GREY, align=PP_ALIGN.CENTER, italic=True)
        if i < len(layers) - 1:
            add_arrow(s, x + w, Inches(3.15), x + w + Inches(0.3), Inches(3.15), color=BLUE)
        x += w + Inches(0.3)
    add_text(s, Inches(0.9), Inches(4.4), Inches(11.5), Inches(0.5),
              "Qt window only wires pages to controllers — Automation / Navigation / Mode / Operations", size=15, color=GREY, align=PP_ALIGN.CENTER)
    add_card(s, Inches(2.0), Inches(5.1), Inches(9.3), Inches(1.4), fill=WHITE)
    add_text(s, Inches(2.3), Inches(5.3), Inches(8.7), Inches(1.0),
              "Pages (view)  ↔  Controllers  ↔  Services\n\"thin window, real logic lives below it\"",
              size=15, color=NAVY, align=PP_ALIGN.CENTER, line_spacing=1.4)

    # ── 4b. Software development practices ────────────────────────────────────
    s = blank_slide(prs)
    add_header(s, "How it was built", "Software engineering, not just a script")
    practices = [
        ("\U0001F501", "Iterative work packages", "WP1–WP7, each with its\nown deliverable"),
        ("\U0001F9EA", "Test-first regression", "184 tests gate every\nbehavioural change"),
        ("\U0001F5A5", "Live VM validation", "Windows + Linux Mint VM,\nnot just CI"),
        ("\U0001F4E6", "Version control", "Git, small reviewable\ncommits per fix"),
    ]
    icon_label_row(s, practices, top=Inches(2.1), circle_d=Inches(1.1), icon_size=30, label_size=15, sub_size=12.5, bg=AMBER)
    add_card(s, Inches(1.2), Inches(4.7), Inches(10.9), Inches(1.6), fill=WHITE)
    add_text(s, Inches(1.5), Inches(4.9), Inches(10.3), Inches(0.4), "17,426 LOC (src) · 3,595 LOC (tests) · 88 Python files · MIT licensed", size=14, bold=True, color=DARK_BLUE)
    add_text(s, Inches(1.5), Inches(5.35), Inches(10.3), Inches(0.85),
              "This session alone: 4 real bugs found and fixed via live VM testing,\na full undo/reset feature added, and cross-platform packaging wired up.",
              size=14, color=GREY, line_spacing=1.25)

    # ── 5. Three modes ────────────────────────────────────────────────────────
    s = blank_slide(prs)
    add_header(s, "User experience", "Same pipeline, three levels of control")
    modes = [
        ("Guided", "\U0001F7E2", GREEN, GREEN_BG, "Zero decisions\nafter mode pick"),
        ("Balanced", "\U0001F535", DARK_BLUE, LIGHT_BLUE_BG, "+ file type\nselection"),
        ("Expert", "\U0001F7E0", AMBER, AMBER_BG, "+ overrides +\nonline verification"),
    ]
    x = Inches(0.9)
    w = Inches(3.6)
    for label, glyph, fg, bg, sub in modes:
        add_card(s, x, Inches(2.0), w, Inches(3.4), fill=WHITE)
        add_icon_circle(s, x + w / 2, Inches(2.65), Inches(0.9), glyph, fg=WHITE, bg=fg, size=24)
        add_text(s, x, Inches(3.25), w, Inches(0.5), label, size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.3), Inches(3.8), w - Inches(0.6), Inches(1.4), sub, size=15, color=GREY,
                  align=PP_ALIGN.CENTER, line_spacing=1.3)
        x += w + Inches(0.2)
    add_card(s, Inches(0.9), Inches(5.65), Inches(11.5), Inches(0.95), fill=AMBER_BG, line_color=AMBER_BG)
    add_text(s, Inches(1.15), Inches(5.83), Inches(11.0), Inches(0.6),
              "Only Expert mode ever touches the network — and only to verify a package\nthat's already been chosen, never to find it.",
              size=14, color=NAVY, line_spacing=1.2)

    # ── 6. Recommendation engine ───────────────────────────────────────────────
    s = blank_slide(prs)
    add_header(s, "Core logic", "Windows app → Linux package, in 4 steps")
    steps = [
        ("\U0001F4CB", "CSV mapping", "238 curated entries"),
        ("\U0001F50D", "Fuzzy match", "handles version numbers"),
        ("\U0001F6E1", "Confidence floor", "never downgrades a\ncurated match"),
        ("\U0001F310", "Repology check", "Expert mode only"),
    ]
    icon_label_row(s, steps, top=Inches(2.1), circle_d=Inches(1.1), icon_size=32, label_size=16, sub_size=13, bg=DARK_BLUE)
    for i in range(len(steps) - 1):
        slot = Inches(11.5) / len(steps)
        cx = Inches(0.9) + slot * i + slot
        add_arrow(s, cx - Inches(0.35), Inches(2.65), cx + Inches(0.05), Inches(2.65), color=BLUE, weight=2.5)
    pic, w, h = add_picture_framed(s, ASSETS / "review_page.png", Inches(3.9), Inches(4.55), Inches(5.6), Inches(2.3),
                                     caption="Real app screen — Review & Confirm")

    # ── 6b. From stated limitation to contribution ─────────────────────────────
    s = blank_slide(prs)
    add_header(s, "Closing a named gap", "This is the term paper's own future work, built")
    add_card(s, Inches(0.7), Inches(1.85), Inches(11.7), Inches(1.9), fill=AMBER_BG, line_color=AMBER_BG)
    add_text(s, Inches(1.0), Inches(2.05), Inches(11.1), Inches(0.4), "THE ORIGINAL PAPER SAID", size=13, bold=True, color=AMBER)
    add_text(s, Inches(1.0), Inches(2.45), Inches(11.1), Inches(0.55),
              "“Compatibility mapping relies on heuristic matching... not guaranteeing feature parity.”", size=16, italic=True, color=NAVY)
    add_text(s, Inches(1.0), Inches(3.05), Inches(11.1), Inches(0.55),
              "Future work named: “moving beyond hard-coded mappings... feature-based matching.”", size=16, italic=True, color=NAVY)
    add_arrow(s, Inches(6.0), Inches(3.95), Inches(6.0), Inches(4.25), color=GREEN, weight=3)
    add_card(s, Inches(0.7), Inches(4.35), Inches(11.7), Inches(2.2), fill=GREEN_BG, line_color=GREEN_BG)
    add_text(s, Inches(1.0), Inches(4.55), Inches(11.1), Inches(0.4), "THIS PRAKTIKUM DELIVERS", size=13, bold=True, color=GREEN)
    deliv = ["Fuzzy matching (SequenceMatcher) replaces exact-string-only lookup",
             "Confidence floors stop the algorithm silently downgrading curated matches",
             "Live Repology verification — package existence checked, not assumed"]
    y = Inches(4.95)
    for d in deliv:
        add_text(s, Inches(1.0), y, Inches(11.1), Inches(0.45), "✅  " + d, size=15, color=NAVY)
        y += Inches(0.48)

    # ── 7. Windows-side workflow ───────────────────────────────────────────────
    s = blank_slide(prs)
    add_header(s, "Workflow — Windows side", "7 guided steps from scan to bundle")
    win_steps = ["Welcome", "Mode", "Scan +\nApp Match", "Settings", "Files", "Review", "Bundle"]
    x = Inches(0.5)
    w = Inches(1.62)
    for i, label in enumerate(win_steps):
        add_icon_circle(s, x + w / 2, Inches(2.2), Inches(0.7), str(i + 1), fg=WHITE, bg=BLUE, size=18)
        add_text(s, x, Inches(2.65), w, Inches(0.6), label, size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        if i < len(win_steps) - 1:
            add_arrow(s, x + w - Inches(0.05), Inches(2.2), x + w + Inches(0.12), Inches(2.2), color=BLUE)
        x += w + Inches(0.12)
    pic, w2, h2 = add_picture_framed(s, ASSETS / "scan_page.png", Inches(4.1), Inches(3.4), Inches(5.4), Inches(3.7),
                                       caption="Real app screen — Scan & App Matching")
    add_text(s, Inches(0.5), Inches(3.6), Inches(3.2), Inches(2.5),
              "Output:\nmigration_bundle.zip\n\nManifest + SHA-256\nApps to install\nSettings & shortcuts\n(optional) Linux binary",
              size=14, color=GREY, line_spacing=1.3)

    # ── 8. Linux-side workflow ─────────────────────────────────────────────────
    s = blank_slide(prs)
    add_header(s, "Workflow — Linux side", "One click: restore → verify → report")
    lin_steps = ["Extract\n& Restore", "Verify\n(SHA-256)", "Report"]
    x = Inches(0.6)
    w = Inches(2.0)
    for i, label in enumerate(lin_steps):
        add_icon_circle(s, x + w / 2, Inches(2.1), Inches(0.85), str(i + 1), fg=WHITE, bg=GREEN, size=22)
        add_text(s, x, Inches(2.6), w, Inches(0.7), label, size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER, line_spacing=1.1)
        if i < len(lin_steps) - 1:
            add_arrow(s, x + w - Inches(0.1), Inches(2.1), x + w + Inches(0.25), Inches(2.1), color=GREEN)
        x += w + Inches(0.25)
    pic, w2, h2 = add_picture_framed(s, ASSETS / "final_report.png", Inches(6.7), Inches(1.9), Inches(5.9), Inches(5.0),
                                       caption="Real report — actual VM restore, 3,344 files")
    add_text(s, Inches(0.6), Inches(3.4), Inches(5.6), Inches(1.0),
              "Failed step? “Restart” or\n“Review & Complete Anyway”\n— nothing fails silently.",
              size=15, color=GREY, line_spacing=1.3)
    add_text(s, Inches(0.6), Inches(4.6), Inches(5.6), Inches(1.6),
              "Reset can fully undo a restore:\nfiles, shortcuts, settings,\nand (opt-in) installed apps.",
              size=15, color=GREY, line_spacing=1.3)

    # ── 9. Privacy ─────────────────────────────────────────────────────────────
    s = blank_slide(prs)
    add_header(s, "Design principle", "Privacy by default, every layer")
    items = [
        ("\U0001F4C1", "Files", "never leave\nthe machine"),
        ("\U0001F310", "Repology", "name/version\nonly"),
        ("\U0001F916", "AI rank", "opt-in,\nExpert only"),
        ("\U0001F4CA", "Metrics", "machine ID\nanonymised"),
        ("\U0001F4E6", "Bundle", "stays local,\nuser controls it"),
    ]
    icon_label_row(s, items, top=Inches(2.3), circle_d=Inches(1.0), icon_size=28, label_size=15, sub_size=12.5, bg=GREEN)

    # ── 10. Validation & reporting ────────────────────────────────────────────
    s = blank_slide(prs)
    add_header(s, "Evidence it worked", "The Sovereignty Score, hash-verified")
    pic, w, h = add_picture_framed(s, ASSETS / "final_report.png", Inches(0.7), Inches(1.85), Inches(5.2), Inches(5.1),
                                     caption="Real report from the actual restore run")
    add_text(s, Inches(6.3), Inches(1.95), Inches(6.2), Inches(0.4), "SOVEREIGNTY SCORE", size=14, bold=True, color=BLUE)
    add_text(s, Inches(6.3), Inches(2.3), Inches(6.2), Inches(0.55), "integrity_score + openness_bonus", size=16, color=NAVY)
    add_text(s, Inches(6.3), Inches(2.75), Inches(6.2), Inches(0.5), "every file hash-verified against the manifest", size=13, italic=True, color=GREY)
    add_text(s, Inches(6.3), Inches(3.6), Inches(6.2), Inches(0.4), "THREE FORMATS", size=14, bold=True, color=BLUE)
    add_text(s, Inches(6.3), Inches(3.95), Inches(6.2), Inches(0.5), "JSON · Markdown · standalone HTML", size=16, color=NAVY)
    add_text(s, Inches(6.3), Inches(4.8), Inches(6.2), Inches(0.4), "NOT JUST PASS/FAIL", size=14, bold=True, color=BLUE)
    add_text(s, Inches(6.3), Inches(5.15), Inches(6.2), Inches(1.1),
              "Warnings, settings applied vs.\nmanual, apps installed vs. failed",
              size=16, color=NAVY, line_spacing=1.2)

    # ── 11. Testing ──────────────────────────────────────────────────────────
    s = blank_slide(prs)
    add_header(s, "Quality", "184 tests, counted directly from the repo")
    add_bar_chart(
        s, Inches(0.8), Inches(1.85), Inches(7.4), Inches(4.3),
        categories=["Unit", "Integration", "Performance", "E2E", "UI"],
        series_name="Test functions",
        values=[72, 68, 18, 15, 11],
        bar_color=DARK_BLUE,
    )
    add_card(s, Inches(8.5), Inches(1.95), Inches(4.0), Inches(4.2), fill=AMBER_BG, line_color=AMBER_BG)
    add_text(s, Inches(8.8), Inches(2.15), Inches(3.4), Inches(0.4), "WHY IT STILL WASN'T ENOUGH", size=12.5, bold=True, color=AMBER)
    add_text(s, Inches(8.8), Inches(2.6), Inches(3.4), Inches(3.3),
              "CI passing ≠ the real app\nworks.\n\nSeveral bugs this session\nwere found only by running\nit live on a Windows +\nLinux Mint VM pair.",
              size=14.5, color=NAVY, line_spacing=1.3)

    # ── 12. Bugs found via live testing ───────────────────────────────────────
    s = blank_slide(prs)
    add_header(s, "Found via live testing", "Real bugs a live VM run exposed")
    bugs = [
        ("One bad package name\nfailed ALL installs", "Batch → falls back to\none-at-a-time"),
        ("Shortcuts matched\nrequested, not installed", "Now matched against\nactual installs only"),
        ("Silent multi-minute\nlog gaps", "Periodic progress\nlogging added"),
        ("Unanswered prompt\ncould hang forever", "120s timeout per\nattempt"),
    ]
    x = Inches(0.6)
    w = Inches(2.95)
    for problem, fix in bugs:
        add_card(s, x, Inches(2.0), w, Inches(1.7), fill=AMBER_BG, line_color=AMBER_BG)
        add_text(s, x + Inches(0.2), Inches(2.15), w - Inches(0.4), Inches(1.3), "⚠  " + problem, size=13, color=NAVY, line_spacing=1.15)
        add_arrow(s, x + w / 2, Inches(3.7), x + w / 2, Inches(3.95), color=GREEN, weight=2.5)
        add_card(s, x, Inches(4.0), w, Inches(1.3), fill=GREEN_BG, line_color=GREEN_BG)
        add_text(s, x + Inches(0.2), Inches(4.15), w - Inches(0.4), Inches(1.05), "✅  " + fix, size=13, color=NAVY, line_spacing=1.15)
        x += w + Inches(0.15)

    # ── 13. Reset feature ──────────────────────────────────────────────────────
    s = blank_slide(prs)
    add_header(s, "New capability", "Reset now undoes everything")
    items = [
        ("\U0001F4C1", "Files"),
        ("\U0001F517", "Shortcuts"),
        ("\U0001F5BC", "Wallpaper"),
        ("\U0001F4E6", "Apps (opt-in)"),
    ]
    x = Inches(1.3)
    for glyph, label in items:
        add_icon_circle(s, x, Inches(2.7), Inches(1.1), glyph, fg=WHITE, bg=GREEN, size=30)
        add_text(s, x - Inches(1.0), Inches(3.35), Inches(2.0), Inches(0.5), label, size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        x += Inches(2.7)
    add_card(s, Inches(1.0), Inches(4.3), Inches(11.3), Inches(2.2), fill=WHITE)
    add_bullet = lambda t, y: add_text(s, Inches(1.35), y, Inches(10.6), Inches(0.5), "•  " + t, size=15, color=NAVY)
    add_bullet("Apps removal is opt-in — off by default, since shared dependencies could be affected", Inches(4.5))
    add_bullet("Works even from a brand-new app session — reads entirely from the saved report", Inches(5.0))
    add_bullet("Live progress in the Activity Log — not just the log file", Inches(5.5))

    # ── 14. Packaging ─────────────────────────────────────────────────────────
    s = blank_slide(prs)
    add_header(s, "Distribution", "Double-click standalone, both platforms")
    add_card(s, Inches(0.9), Inches(2.0), Inches(5.4), Inches(3.6), fill=WHITE)
    add_icon_circle(s, Inches(3.6), Inches(2.85), Inches(1.0), "\U0001FA9F", fg=WHITE, bg=DARK_BLUE, size=30)
    add_text(s, Inches(1.1), Inches(3.55), Inches(5.0), Inches(0.5), "MigrationWizard.exe", size=17, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.3), Inches(4.05), Inches(4.6), Inches(1.4),
              "PyInstaller build of the full Qt app.\nCan bake the Linux binary in directly\nfor a true single-file standalone.",
              size=13.5, color=GREY, align=PP_ALIGN.CENTER, line_spacing=1.3)
    add_card(s, Inches(6.9), Inches(2.0), Inches(5.4), Inches(3.6), fill=WHITE)
    add_icon_circle(s, Inches(9.6), Inches(2.85), Inches(1.0), "\U0001F427", fg=WHITE, bg=GREEN, size=30)
    add_text(s, Inches(7.1), Inches(3.55), Inches(5.0), Inches(0.5), "MigrationWizard (ELF)", size=17, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(7.3), Inches(4.05), Inches(4.6), Inches(1.4),
              "Built on the target distro itself.\nEmbedded automatically into every\nbundle the Windows app creates.",
              size=13.5, color=GREY, align=PP_ALIGN.CENTER, line_spacing=1.3)
    add_text(s, Inches(0.9), Inches(5.9), Inches(11.4), Inches(0.6),
              "Bundle stays self-contained: unzip → run → restore. No Python install needed on either side.",
              size=14, italic=True, color=GREY, align=PP_ALIGN.CENTER)

    # ── 14b. Project management — work packages & timeline ─────────────────────
    s = blank_slide(prs)
    add_header(s, "Project management", "Seven work packages, April → June")
    wps = [
        ("WP1", "Architecture", "Apr 1–17", GREEN),
        ("WP2", "Recommendations", "Apr 15 – May 14", GREEN),
        ("WP3", "Workflow & UI", "Apr 15 – May 27", GREEN),
        ("WP4", "Backup Pipeline", "May 1–27", GREEN),
        ("WP5", "Linux Restore", "May 10–28", GREEN),
        ("WP6", "Packaging", "May 25–29", GREEN),
        ("WP7", "Quality, Docs & Presentation", "Apr 1 – Jun 8", AMBER),
    ]
    y = Inches(1.85)
    row_h = Inches(0.62)
    for code, name, dates, color in wps:
        add_card(s, Inches(0.7), y, Inches(1.0), row_h - Pt(6), fill=color, line_color=color)
        add_text(s, Inches(0.7), y + Pt(4), Inches(1.0), row_h - Pt(6), code, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(1.85), y + Pt(2), Inches(4.6), row_h, name, size=14.5, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(6.6), y + Pt(2), Inches(2.8), row_h, dates, size=13, color=GREY, anchor=MSO_ANCHOR.MIDDLE)
        status = "Done" if color == GREEN else "Active"
        add_pill(s, Inches(9.6), y + Pt(6), Inches(1.4), row_h - Pt(14), status, WHITE, color, size=12)
        y += row_h + Pt(4)
    add_text(s, Inches(0.7), y + Pt(8), Inches(11.5), Inches(0.5),
              "Risk buffer (May 29 – Jun 4) reserved for fixes and real-environment testing before this presentation.",
              size=13.5, italic=True, color=GREY)

    # ── 14c. Project management — risk register & success criteria ─────────────
    s = blank_slide(prs)
    add_header(s, "Project management", "Risk register & success criteria")
    risks = [
        ("PySide6 packaging fails", "High impact", "Mitigated", GREEN),
        ("Files restored to wrong location", "High impact", "Fixed", GREEN),
        ("App recs miss niche software", "Medium impact", "Mitigated", GREEN),
        ("No E2E test on real hardware", "Medium impact", "Open", AMBER),
    ]
    add_text(s, Inches(0.7), Inches(1.8), Inches(5.6), Inches(0.4), "KEY RISKS", size=13, bold=True, color=BLUE)
    y = Inches(2.2)
    for risk, impact, status, color in risks:
        add_card(s, Inches(0.7), y, Inches(5.6), Inches(0.85), fill=WHITE)
        add_text(s, Inches(0.9), y + Inches(0.08), Inches(3.6), Inches(0.4), risk, size=13.5, bold=True, color=NAVY)
        add_text(s, Inches(0.9), y + Inches(0.45), Inches(2.0), Inches(0.35), impact, size=11.5, color=GREY)
        add_pill(s, Inches(4.7), y + Inches(0.22), Inches(1.4), Inches(0.42), status, WHITE, color, size=11)
        y += Inches(0.98)
    add_text(s, Inches(6.7), Inches(1.8), Inches(5.9), Inches(0.4), "SUCCESS CRITERIA", size=13, bold=True, color=BLUE)
    crit = [
        ("Full scan → backup → restore pipeline", "Done"),
        ("80% of common apps mapped", "238 mapped"),
        ("150+ automated tests", "184 passing"),
        ("Single double-click package", "exe + ELF"),
    ]
    y = Inches(2.2)
    for label, result in crit:
        add_card(s, Inches(6.7), y, Inches(5.9), Inches(0.85), fill=GREEN_BG, line_color=GREEN_BG)
        add_text(s, Inches(6.9), y + Inches(0.08), Inches(4.3), Inches(0.4), label, size=13.5, color=NAVY)
        add_text(s, Inches(6.9), y + Inches(0.45), Inches(4.3), Inches(0.35), "✅ " + result, size=12.5, bold=True, color=GREEN)
        y += Inches(0.98)

    # ── 15. Evaluation ─────────────────────────────────────────────────────────
    s = blank_slide(prs)
    add_header(s, "Evaluation", "What the project delivers")
    metrics = [("238", "app mappings"), ("184", "automated tests"), ("3", "guidance modes"), ("2", "platforms packaged")]
    x = Inches(0.8)
    for num, label in metrics:
        add_card(s, x, Inches(1.85), Inches(2.75), Inches(1.7), fill=WHITE)
        add_text(s, x, Inches(2.05), Inches(2.75), Inches(0.8), num, size=36, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.15), Inches(2.8), Inches(2.45), Inches(0.6), label, size=13, color=NAVY, align=PP_ALIGN.CENTER)
        x += Inches(2.88)
    add_bar_chart(
        s, Inches(0.8), Inches(3.75), Inches(5.5), Inches(3.3),
        categories=["Before this session", "Now"], series_name="App mappings",
        values=[149, 238], bar_color=GREEN, title="App mapping table growth",
    )
    pic, w, h = add_picture_framed(s, ASSETS / "final_report.png", Inches(6.7), Inches(3.85), Inches(5.9), Inches(3.0),
                                     caption="100% sovereignty score — real restore, 3,344 files")

    # ── 16. Conclusion & future work ──────────────────────────────────────────
    s = blank_slide(prs)
    add_header(s, "Wrap-up", "Conclusion & future work")
    add_card(s, Inches(0.7), Inches(1.9), Inches(5.7), Inches(4.7), fill=GREEN_BG, line_color=GREEN_BG)
    add_text(s, Inches(1.0), Inches(2.1), Inches(5.0), Inches(0.4), "✅  DELIVERED", size=15, bold=True, color=GREEN)
    deliv = ["Dynamic recommendation engine", "Unified mode policy (Qt + CLI)", "Per-stage execution timing",
             "184 test regression coverage", "Real-VM-tested reliability fixes", "Full undo/reset + packaging"]
    y = Inches(2.65)
    for d in deliv:
        add_text(s, Inches(1.1), y, Inches(5.1), Inches(0.45), "•  " + d, size=14.5, color=NAVY)
        y += Inches(0.62)
    add_card(s, Inches(6.7), Inches(1.9), Inches(5.7), Inches(4.7), fill=AMBER_BG, line_color=AMBER_BG)
    add_text(s, Inches(7.0), Inches(2.1), Inches(5.0), Inches(0.4), "→  NEXT", size=15, bold=True, color=AMBER)
    nxt = ["Real E2E test, physical machines", "Live USB write automation", "Resume/checkpoint for large restores", "Broader distro support"]
    y = Inches(2.65)
    for n in nxt:
        add_text(s, Inches(7.1), y, Inches(5.1), Inches(0.6), "•  " + n, size=14.5, color=NAVY)
        y += Inches(0.7)

    # ── 16b. References ─────────────────────────────────────────────────────────
    s = blank_slide(prs)
    add_header(s, "References", "Sovereignty framing cited from the term paper")
    refs = [
        "Floridi, L. (2020). The fight for digital sovereignty. Philosophy & Technology, 33(3), 369–378.",
        "Pohle, J., & Thiel, T. (2020). Digital sovereignty. Internet Policy Review, 9(4).",
        "Roberts, H., et al. (2021). Safeguarding European values with digital sovereignty. Internet Policy Review, 10(3).",
        "Bechara, J., & Lechner, U. (2024). Digital sovereignty and open-source software. I4CS 2024, Springer.",
        "Wehnes, H. (2024). Preventing digital colony and lock-in. INFORMATIK 2024, Gesellschaft für Informatik.",
        "European Commission (2020, 2022). Shaping Europe's digital future / A digital decade for Europe.",
        "Wire (2025). The state of digital sovereignty in Europe.",
        "StatCounter Global Stats (2025). Desktop OS Market Share, Europe.",
    ]
    y = Inches(1.95)
    for r in refs:
        add_text(s, Inches(0.9), y, Inches(11.5), Inches(0.45), "•  " + r, size=14, color=NAVY)
        y += Inches(0.58)
    add_text(s, Inches(0.9), y + Inches(0.1), Inches(11.5), Inches(0.5),
              "Full bibliography in: Arthur, J. — Digital Sovereignty Through Semi-Automated OS Migration (term paper).",
              size=13, italic=True, color=GREY)

    # ── 17. Thank you ──────────────────────────────────────────────────────────
    s = blank_slide(prs, fill=NAVY)
    add_icon_circle(s, Inches(2.3), Inches(1.5), Inches(1.0), "\U0001FA9F", fg=WHITE, bg=DARK_BLUE, size=36)
    add_text(s, Inches(2.85), Inches(1.05), Inches(0.9), Inches(0.9), "→", size=36, bold=True, color=BLUE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_icon_circle(s, Inches(4.0), Inches(1.5), Inches(1.0), "\U0001F427", fg=WHITE, bg=GREEN, size=36)
    add_text(s, Inches(0.8), Inches(2.9), Inches(11.5), Inches(1.2), "Thank you", size=48, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(3.9), Inches(11.5), Inches(0.6), "Questions & live demo", size=22, color=LIGHT_BLUE_BG)

    return prs


if __name__ == "__main__":
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = build()
    prs.save(OUT_PATH)
    print(f"Saved {len(prs.slides)} slides to {OUT_PATH}")
